import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def get_image_path(pattern):
    base_dir = r"C:\Users\UDAYV\.gemini\antigravity\brain\1e9a45dd-5dac-41ac-bdd9-3b16f65ee7fd"
    search_path = os.path.join(base_dir, pattern)
    files = glob.glob(search_path)
    if files:
        # Get the latest one if multiple exist
        files.sort(key=os.path.getmtime)
        return files[-1]
    return None

def clear_and_write_paragraphs(tf, paragraphs, font_size=12, custom_spacing=1.15):
    # Clear existing paragraphs by modifying the first and deleting the rest
    p0 = tf.paragraphs[0]
    p0.text = ""
    
    # Remove all paragraphs except the first one
    while len(tf.paragraphs) > 1:
        p_to_remove = tf.paragraphs[1]
        tf._element.remove(p_to_remove._element)
        
    tf.word_wrap = True
    
    for idx, text in enumerate(paragraphs):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = text
        p.space_after = Pt(4)
        p.line_spacing = custom_spacing
        
        # Set fonts
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Professional charcoal/dark gray
        
        # Check formatting flags
        if text.startswith("### "):
            # Subheader
            p.text = text.replace("### ", "")
            p.font.bold = True
            p.font.size = Pt(font_size + 3)
            p.font.color.rgb = RGBColor(0x0B, 0x25, 0x45) # Deep Navy Blue
            p.space_before = Pt(8)
        elif text.startswith("**") and text.endswith("**"):
            # Main Header
            p.text = text.replace("**", "")
            p.font.bold = True
            p.font.size = Pt(font_size + 6)
            p.font.color.rgb = RGBColor(0x0B, 0x25, 0x45) # Deep Navy Blue
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        elif text.startswith("- **") or text.startswith("* **"):
            # Bullet point with bold prefix
            p.font.size = Pt(font_size)
            cleaned = text[4:]
            bold_part, regular_part = cleaned.split("**", 1) if "**" in cleaned else (cleaned, "")
            p.text = ""
            run1 = p.add_run()
            run1.text = "•  " + bold_part
            run1.font.bold = True
            run1.font.name = "Arial"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            
            if regular_part:
                run2 = p.add_run()
                run2.text = regular_part
                run2.font.bold = False
                run2.font.name = "Arial"
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        elif text.startswith("- ") or text.startswith("* "):
            p.text = "•  " + text[2:]
            p.font.bold = False
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        elif text.startswith("  "):
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

def main():
    ppt_path = r"d:\BUNNY\PROJECTS\ISRO\Idea_Submission_Template.pptx"
    output_path = r"d:\BUNNY\PROJECTS\ISRO\VAYUSETU_ISRO_BAH_2026_Idea_Submission.pptx"
    logo_path = r"d:\BUNNY\PROJECTS\ISRO\frontend\public\logo\logo.jpg"
    
    prs = Presentation(ppt_path)
    print("Successfully opened the template!")
    
    # ------------------ SLIDE 1: TITLE ------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Team Name" in text:
                shape.text_frame.text = "Team Name : ClimateX Labs"
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(20)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0B, 0x25, 0x45)
            elif "Team Leader Name" in text:
                shape.text_frame.text = "Team Leader Name : Kalle Uday Bhaskar"
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(18)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0B, 0x25, 0x45)
            elif "Problem Statement" in text:
                shape.text_frame.text = "Problem Statement : Challenge 5 - AI-Powered Climate Digital Twin of India (AP Pilot Scope)"
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(16)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Insert logo on Slide 1
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(6.6), Inches(0.4), width=Inches(2.8))
        print("Inserted VAYUSETU Logo on Slide 1 successfully!")

    # ------------------ SLIDE 2: TEAM MEMBERS ------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_table:
            table = shape.table
            cell = table.cell(0, 0)
            cell.text = "Team Leader:\n\nName: Kalle Uday Bhaskar\nRole: Geospatial & Software Architecture\nCollege: Indian Institute of Space Science and Technology"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"
            
            cell = table.cell(0, 1)
            cell.text = "Team Member-1:\n\nName: [Name]\nRole: Deep Learning Engineer (XAI & ConvLSTM)\nCollege: Geospatial Science Division"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"
                
            cell = table.cell(1, 0)
            cell.text = "Team Member-2:\n\nName: [Name]\nRole: Mobile-Responsive Frontend & GIS Developer\nCollege: Earth & Climate Applications Division"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"
                
            cell = table.cell(1, 1)
            cell.text = "Team Member-3:\n\nName: [Name]\nRole: Data Engineer (Ingestion & Validation)\nCollege: Space Applications Centre"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"

    # ------------------ SLIDE 3: OPPORTUNITY & USP ------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame and "Opportunity should be able to" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.6)
            shape.width = Inches(9.0)
            shape.height = Inches(4.7)
            
            paragraphs = [
                "**VAYUSETU: OPPORTUNITY, PILOT SCOPING & USP**",
                "### Scoped Pilot Region: Coastal Andhra Pradesh",
                "- **High Climate Variability**: Vulnerable to cyclones, flood catchment discharges, and severe heatwaves. Excellent PoC region for rainfall/temperature forecasting and early warnings.",
                "- **The Gap**: Traditional models are slow. Visual observation portals (e.g. standard dashboards) lack real-time digital twin scenarios and explainable prediction mechanics.",
                "### How VAYUSETU Solves the Problem",
                "- **Indigenous Data Fusion**: Ingests INSAT-3D LST, SST, and IMD Pune 0.25° gridded datasets into a PostGIS coordinate grid.",
                "- **AI-Accelerated Forecasts**: Replaces heavy physical model equations with fast ConvLSTM models to deliver spatial projections in seconds.",
                "### Unique Selling Proposition (USP)",
                "- **USP #1**: Scoped pilot digital twin using indigenous INSAT and IMD inputs.",
                "- **USP #2**: What-If Scenario Builder (e.g. +20% rain anomaly) simulating subsequent flood catchment and heat shifts.",
                "- **USP #3**: Explainable AI (XAI) outputting clear contributing factors for public trust."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)

    # ------------------ SLIDE 4: FEATURES OFFERED ------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame and "List of features offered" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.6)
            shape.width = Inches(9.0)
            shape.height = Inches(4.7)
            
            paragraphs = [
                "**CORE FUNCTIONALITIES OF VAYUSETU**",
                "### 1. Flood Early Warning Module (FEWS)",
                "- **Precipitation Anomalies**: Input forecast precipitation (e.g. 220mm) to output dynamic flood probabilities, target drainage basins, and NDRF advisories.",
                "### 2. Explainable AI (XAI) Contributor Panel",
                "- **Feature Attribution**: Explains rainfall forecasts using SHAP values (e.g. SST anomaly contribution: 34%, Humidity: 28%, Wind Vectors: 21%), avoiding black-box distrust.",
                "### 3. Climate Risk Score Indexing",
                "- **Aggregated Scores**: Calculates dynamic indices for Flood, Heatwave, and Drought, yielding an overall regional alert level.",
                "### 4. Interactive What-If Scenario Builder",
                "- **Variables**: Adjust Precipitation, Temp Rise, and Urban expansion sliders. Instantly outputs projected risk shifts (e.g. Urbanization +15% -> Flood Risk +34%).",
                "### 5. Timeline Slider & Model Quality",
                "- **Temporal Playback**: Past State -> Current Ingestion -> 24h -> 48h Forecast -> Scenario projection.",
                "- **Production Monitoring**: Real-time dashboards monitoring prediction accuracy (92%) and data drift (1.8%)."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)

    # ------------------ SLIDE 5: PROCESS FLOW ------------------
    slide5 = prs.slides[4]
    img_path_5 = get_image_path("gov_process_flow_*.png")
    for shape in slide5.shapes:
        if shape.has_text_frame and "Process flow diagram" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.5)
            shape.width = Inches(9.0)
            shape.height = Inches(0.7)
            
            paragraphs = [
                "**VAYUSETU PROCESS FLOW: PIPELINE REPRODUCIBILITY**",
                "Clean ingestion of satellite feeds, data harmonizing, AI forecasting, digital twin updates, and decision advisories."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)
            
    if img_path_5:
        slide5.shapes.add_picture(img_path_5, Inches(1.2), Inches(1.3), width=Inches(7.6))
        print("Inserted Slide 5 Process Flow Diagram successfully!")

    # ------------------ SLIDE 6: WIREFRAME ------------------
    slide6 = prs.slides[5]
    img_path_6 = get_image_path("gov_dashboard_wireframe_*.png")
    for shape in slide6.shapes:
        if shape.has_text_frame and "Wireframes/Mock diagrams" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.5)
            shape.width = Inches(9.0)
            shape.height = Inches(0.7)
            
            paragraphs = [
                "**VAYUSETU GOVERNMENT-GRADE OPERATIONS INTERFACE**",
                "Minimalist layout featuring spatial maps, scenario sliders, timeline playbacks, risk tables, and RAG advisory chatbot."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)
            
    if img_path_6:
        slide6.shapes.add_picture(img_path_6, Inches(1.2), Inches(1.3), width=Inches(7.6))
        print("Inserted Slide 6 Dashboard Wireframe successfully!")

    # ------------------ SLIDE 7: ARCHITECTURE ------------------
    slide7 = prs.slides[6]
    img_path_7 = get_image_path("gov_architecture_diagram_*.png")
    for shape in slide7.shapes:
        if shape.has_text_frame and "Architecture diagram of the" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.5)
            shape.width = Inches(9.0)
            shape.height = Inches(0.7)
            
            paragraphs = [
                "**VAYUSETU DECOUPLED TECHNICAL ARCHITECTURE**",
                "3-tier layout: Mobile-responsive Next.js UI, FastAPI microservice router, and PostgreSQL core GIS engine."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)
            
    if img_path_7:
        slide7.shapes.add_picture(img_path_7, Inches(1.2), Inches(1.3), width=Inches(7.6))
        print("Inserted Slide 7 Architecture Diagram successfully!")

    # ------------------ SLIDE 8: TECHNOLOGIES ------------------
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.has_text_frame and "Technologies to be used" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.6)
            shape.width = Inches(9.0)
            shape.height = Inches(4.7)
            
            paragraphs = [
                "**VAYUSETU PRODUCTION-GRADE TECH STACK**",
                "### 1. Presentation Layer (Government-Grade Minimalism)",
                "- **Next.js 15 & React**: Type-safe responsive presentation featuring ShadCN UI widgets.",
                "- **Tailwind CSS**: Media-responsive viewport scaling ensuring absolute mobile-responsive compatibility.",
                "- **Leaflet.js & Mapbox**: Touch-friendly rendering of GIS coordinate layers.",
                "### 2. Service & API Layer",
                "- **FastAPI**: Decoupled asynchronous REST gateway offering OAuth2 JWT authentication and rate limiting.",
                "- **Redis Cache**: Caching frequent raster data calls to optimize throughput.",
                "### 3. Data & ML Modeling Layer",
                "- **PostgreSQL with PostGIS**: Relational storage for spatial raster coordinates and metadata.",
                "- **PyTorch, TensorFlow, XGBoost**: AI core for Spatio-temporal and temperature forecasting.",
                "- **Rasterio, GeoPandas, GDAL**: Heavy scientific packages to ingest NetCDF/GeoTIFF datasets.",
                "- **SHAP & LIME**: Explainable AI frameworks for feature importance calculations."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)

    # ------------------ SLIDE 9: FEASIBILITY, COST & SCALABILITY ------------------
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.has_text_frame and "Estimated implementation cost" in shape.text_frame.text:
            shape.left = Inches(0.5)
            shape.top = Inches(0.6)
            shape.width = Inches(9.0)
            shape.height = Inches(4.7)
            
            paragraphs = [
                "**FEASIBILITY, SECURITY, COSTS & SCALABILITY**",
                "### 1. Security & Data Governance",
                "- **Authentication**: Secure JWT access and HTTP-only refresh tokens. Role-Based Access Control (RBAC).",
                "- **Encryption**: Local database encryption-at-rest (AES-256) and secure TLS 1.3 transport.",
                "### 2. Operational Feasibility & Costs",
                "- **Open-Source Stack**: Open databases and libraries avoid licensing fees.",
                "- **Cloud Host Estimate (Monthly)**: ~₹8,000 to ₹10,000 for pilot phase compute, PostGIS, and Redis cache.",
                "- **Large-Scale Deploy (NIC/ISRO Cloud)**: ~₹2 - 5 Lakhs annually for multi-node GPU clustering.",
                "### 3. National Scalability Pathway",
                "- **GTM Path**: AP Coastal Region Pilot → State Level Rollout → Basin-Level Expansion → National Climate Twin (ISRO-Bhuvan integration).",
                "### 4. IP Potential",
                "- **Hydrometeorological Runoff Engine**: Spatio-Temporal Climate Fusion and Data Assimilation Engine."
            ]
            clear_and_write_paragraphs(shape.text_frame, paragraphs, font_size=10)

    # ------------------ SLIDE 10: CLOSING QUOTE ------------------
    slide10 = prs.slides[9]
    if os.path.exists(logo_path):
        slide10.shapes.add_picture(logo_path, Inches(3.8), Inches(0.4), width=Inches(2.4))
        print("Inserted VAYUSETU Logo on Slide 10 successfully!")

    tx_box = slide10.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(8.0), Inches(2.2))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = '"A resilient nation is built not by predicting the future, but by preparing for it. VAYUSETU transforms India\'s climate data into actionable intelligence for a safer, smarter, and more sustainable India."'
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0B, 0x25, 0x45)
    p.alignment = 1 # Center
    
    p2 = tf.add_paragraph()
    p2.text = "\n— ClimateX Labs | ISRO BAH 2026"
    p2.font.name = "Arial"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p2.alignment = 1 # Center
    
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
